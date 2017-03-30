# !/usr/bin/python
# -*- coding: utf-8 -*-
from flask import Flask, render_template, request, send_from_directory
from pptx import Presentation
import textslide.summarize as ts
import os

app = Flask(__name__)



@app.route('/Result.html', methods=['POST'])
def summ():
    ss = ts.SimpleSummarizer()
    text = request.form['TokenText']
    n = request.form['sentnum']
    name = request.form['Name']
    split = str(text).split("\n")
    a = []
    for items in split:
        sum1 = ss.summarize(items, n)
        a.append(sum1)
    sss = ts.WordFreq()
    words = sss.top(text)
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = words[0][0]
    subtitle.text = name
    v = (a.__len__()) - 1
    for items in a:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "test"
        subtitle.text = items
        prs.save('test.pptx')
    return send_from_directory("", "test.pptx", as_attachment=True,
                               mimetype='application/octet-stream')



@app.route('/')
def index():
    return render_template('index.html')


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)
